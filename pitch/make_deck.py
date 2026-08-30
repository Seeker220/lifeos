"""Generates the LifeOS pitch deck.

Colors and type are lifted from the app's own dark theme (ui/theme/Color.kt) so
the deck and the product look like the same thing.

    python3 pitch/make_deck.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---- palette (ui/src/main/kotlin/com/lifeos/ui/theme/Color.kt) ----
BACKDROP = RGBColor(0x07, 0x09, 0x0C)
SURFACE0 = RGBColor(0x0B, 0x0E, 0x13)
SURFACE1 = RGBColor(0x11, 0x15, 0x1C)
SURFACE2 = RGBColor(0x17, 0x1C, 0x25)
BORDER = RGBColor(0x1E, 0x24, 0x2F)
ACCENT = RGBColor(0xA8, 0xC7, 0xFA)
ACCENT_VIVID = RGBColor(0x4C, 0x8D, 0xFF)
ACCENT_DEEP = RGBColor(0x28, 0x49, 0x7A)
ACCENT_INK = RGBColor(0x0A, 0x30, 0x5F)
WARN = RGBColor(0xFF, 0xD8, 0xA8)
DANGER = RGBColor(0xFF, 0xB4, 0xAB)
SUCCESS = RGBColor(0x5A, 0xF0, 0xBE)
VIOLET = RGBColor(0xC7, 0xA9, 0xFF)
T1 = RGBColor(0xE8, 0xEE, 0xF5)
T2 = RGBColor(0x97, 0xA3, 0xB2)
T3 = RGBColor(0x61, 0x6C, 0x7A)

DISPLAY = "Montserrat"
BODY = "Noto Sans"

W, H = 13.333, 7.5
M = 0.85
CW = W - 2 * M

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_slide_no = 0


def slide(bg=SURFACE0, numbered=True):
    global _slide_no
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    if numbered:
        _slide_no += 1
        text(s, W - M - 0.8, H - 0.62, 0.8, 0.3, [("%02d" % _slide_no, 9, T3, False)],
             align=PP_ALIGN.RIGHT)
        text(s, M, H - 0.62, 3.0, 0.3, [("LifeOS", 9, T3, False)])
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=1.3, anchor=MSO_ANCHOR.TOP,
         space_after=0):
    """runs: list of (text, size_pt, color, bold) or list of such lists (paragraphs)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if space_after:
            p.space_after = Pt(space_after)
        for item in para:
            body, size, color, bold = item[0], item[1], item[2], item[3]
            font_name = item[4] if len(item) > 4 else None
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font_name or (DISPLAY if size >= 20 or bold else BODY)
    return tb


def card(s, x, y, w, h, fill=SURFACE1, border=BORDER, radius=0.045):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def bar(s, x, y, w, h, color=ACCENT_VIVID):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def head(s, eyebrow, title, title_size=34, sub=None):
    bar(s, M, 0.66, 0.055, 0.2)
    text(s, M + 0.2, 0.62, CW, 0.3, [(eyebrow.upper(), 10.5, ACCENT, True)])
    text(s, M, 1.02, CW - 0.4, 1.0, [(title, title_size, T1, True)], spacing=1.12)
    if sub:
        line_h = 0.62 if title_size >= 30 else 0.52
        lines = title.count("\n") + 1
        text(s, M, 1.02 + lines * line_h, CW - 1.2, 0.6,
             [(sub, 13, T2, False)], spacing=1.35)


def feature_card(s, x, y, w, h, title, body, accent=ACCENT, pad=0.26):
    card(s, x, y, w, h)
    bar(s, x + pad, y + pad, 0.42, 0.045, accent)
    text(s, x + pad, y + pad + 0.22, w - 2 * pad, 0.36, [(title, 14, T1, True)])
    text(s, x + pad, y + pad + 0.66, w - 2 * pad, h - pad * 2 - 0.66,
         [(body, 10.5, T2, False)], spacing=1.32)


def bullets(s, x, y, w, items, size=12, gap=0.42, marker="—", mcolor=ACCENT_VIVID,
            color=T2):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            runs = [(marker + "  ", size, mcolor, False),
                    (lead, size, T1, True),
                    ("  " + rest, size, color, False)]
        else:
            runs = [(marker + "  ", size, mcolor, False), (it, size, color, False)]
        text(s, x, y + i * gap, w, gap, runs, spacing=1.28)


def stat(s, x, y, w, value, label, color=T1):
    text(s, x, y, w, 0.6, [(value, 30, color, True)])
    text(s, x, y + 0.56, w, 0.5, [(label, 10, T2, False)], spacing=1.25)


# ==========================================================================
# 01 — Title
# ==========================================================================
s = slide(bg=BACKDROP, numbered=False)
bar(s, M, 2.35, 0.055, 1.55)
text(s, M + 0.28, 2.28, 9.0, 1.2, [("LifeOS", 62, T1, True)])
text(s, M + 0.32, 3.42, 10.5, 0.6,
     [("Stop managing tasks. Start executing goals.", 21, ACCENT, False)])
text(s, M + 0.32, 4.15, 9.4, 1.1,
     [("A personal AI operating system that doesn't just plan your life — "
       "it enforces it, on the device, whether you feel like it or not.", 13.5, T2, False)],
     spacing=1.42)
text(s, M, 6.35, 6.0, 0.4,
     [("Hackathon prototype  ·  Native Android  ·  Working APK", 11, T3, False)])

# ==========================================================================
# 02 — The hook
# ==========================================================================
s = slide()
head(s, "Where this idea came from", "Ever felt like you're\nprocrastinating too much?", 32)

text(s, M, 2.55, 6.9, 2.4,
     [[("Your exam is in three weeks. The interview is a month out. "
        "You've been meaning to get fit since January.", 14.5, T1, False)],
      [("So why are you forty minutes into a reels binge, or six tabs deep "
        "in a Reddit rabbit hole?", 14.5, T1, False)]],
     spacing=1.45, space_after=12)

text(s, M, 5.05, 6.9, 0.5, [("At least, I do.", 17, ACCENT, True)])
text(s, M, 5.62, 6.9, 0.8,
     [("That's not a market slide. That's my screen time — and probably yours.",
       12, T2, False)], spacing=1.35)

card(s, 8.15, 2.5, 4.3, 3.5, fill=SURFACE1)
text(s, 8.5, 2.85, 3.6, 0.4, [("THE HONEST VERSION", 9.5, T3, True)])
bullets(s, 8.5, 3.4, 3.6,
        ["I don't need another list.",
         "I don't need a nicer dashboard.",
         "I know exactly what I should be doing.",
         "I need something that will actually stop me."],
        size=11.5, gap=0.55)

# ==========================================================================
# 03 — Planning vs enforcement
# ==========================================================================
s = slide()
head(s, "Why every productivity app fails", "You don't lack a plan.\nYou lack enforcement.", 32)

card(s, M, 2.9, 5.35, 3.35)
text(s, M + 0.35, 3.2, 4.6, 0.4, [("WHAT EVERY APP GIVES YOU", 9.5, T3, True)])
bullets(s, M + 0.35, 3.72, 4.6,
        ["Lists, boards, and tags", "Reminders you swipe away",
         "Streaks you abandon on day four", "A chatbot that writes a perfect plan"],
        size=11.5, gap=0.5, mcolor=T3)

card(s, M + 5.65, 2.9, 5.35, 3.35, fill=SURFACE2)
text(s, M + 6.0, 3.2, 4.6, 0.4, [("WHAT ACTUALLY STOPS YOU", 9.5, T3, True)])
text(s, M + 6.0, 3.75, 4.6, 1.0, [("Nothing.", 30, DANGER, True)])
text(s, M + 6.0, 4.62, 4.6, 1.3,
     [("Every one of them is a suggestion. The plan exists, sits there, "
       "and Instagram still wins.", 12, T2, False)], spacing=1.38)

text(s, M, 6.42, CW, 0.5,
     [("People don't fail from disorganization. They fail at ", 13, T2, False),
      ("execution", 13, T1, True),
      (" — and execution is the part nobody built.", 13, T2, False)])

# ==========================================================================
# 04 — Market gap / incentives
# ==========================================================================
s = slide()
head(s, "Why the gap exists", "Nobody who owns your attention\nwill help you escape it.", 30)

cw3 = (CW - 0.56) / 3
feature_card(s, M, 2.95, cw3, 2.5, "Google",
             "YouTube watch time is the revenue line. An app that gets you off "
             "YouTube is an app that cuts their own numbers.", DANGER)
feature_card(s, M + cw3 + 0.28, 2.95, cw3, 2.5, "Meta",
             "Reels engagement is the product. Nothing they ship will ever be "
             "designed to make you scroll less.", DANGER)
feature_card(s, M + 2 * (cw3 + 0.28), 2.95, cw3, 2.5, "OS-level Screen Time",
             "Shipped, but deliberately toothless. One tap to dismiss, because "
             "friction costs them engagement too.", WARN)

card(s, M, 5.72, CW, 1.0, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 0.35, 5.98, CW - 0.7, 0.6,
     [("A third-party OS can be opinionated on purpose. ", 14, T1, True),
      ("That misalignment is the entire opening — and it isn't closing.",
       14, T2, False)])

# ==========================================================================
# 05 — What LifeOS is
# ==========================================================================
s = slide()
head(s, "The product", "One sentence becomes\nenforced device state.", 32)

card(s, M, 2.95, 5.5, 2.35, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 0.35, 3.2, 4.8, 0.3, [("YOU SAY", 9.5, ACCENT, True)])
text(s, M + 0.35, 3.62, 4.8, 1.5,
     [("\u201cI have a DSA assignment due Tuesday, I keep doomscrolling "
       "Instagram, wake me at 7.\u201d", 14, T1, False)], spacing=1.4)

text(s, M + 5.75, 3.85, 0.5, 0.5, [("\u2192", 26, ACCENT_VIVID, True)])

card(s, M + 6.35, 2.55, CW - 6.35, 3.7)
text(s, M + 6.7, 2.83, 4.2, 0.3, [("LIFEOS WRITES", 9.5, ACCENT, True)])
bullets(s, M + 6.7, 3.25, 4.3,
        [("Goal", "with a hard Tuesday deadline"),
         ("Tasks + study blocks", "on your calendar"),
         ("Instagram", "capped, then blocked"),
         ("7:00 alarm", "armed in the system clock"),
         ("Memory", "\u201cfolds to Instagram at night\u201d")],
        size=11.5, gap=0.6, marker="\u2713", mcolor=SUCCESS)

text(s, M, 6.45, CW, 0.5,
     [("It doesn't reply with a checklist. It reaches into the phone and changes it.",
       13.5, T1, True)])

# ==========================================================================
# 06 — How it works
# ==========================================================================
s = slide()
head(s, "How it works", "Intent \u2192 Plan \u2192 State \u2192 Enforce \u2192 Reflect", 30,
     sub="One model call returns a spoken reply plus typed actions. One component applies them.")

steps = [
    ("Intent", "Chat, or an email that just landed", ACCENT),
    ("Plan", "Model returns reply + actions[]", VIOLET),
    ("State", "Durable life state, outside the LLM", ACCENT_VIVID),
    ("Enforce", "Alarms, overlay, timeouts, DNS", DANGER),
    ("Reflect", "Today, Goals, Focus, XP", SUCCESS),
]
bw = (CW - 4 * 0.3) / 5
for i, (t, b, c) in enumerate(steps):
    x = M + i * (bw + 0.3)
    feature_card(s, x, 2.95, bw, 2.2, t, b, c, pad=0.22)
    if i < 4:
        text(s, x + bw + 0.05, 3.85, 0.3, 0.3,
             [("\u203a", 17, T3, True)], align=PP_ALIGN.CENTER)

card(s, M, 5.55, CW, 1.15, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 0.35, 5.82, CW - 0.7, 0.7,
     [("The model is the planner. The database is the soul. ", 14, T1, True),
      ("Exactly one component writes state, so the AI can never corrupt "
       "your goals — only propose changes to them.", 13, T2, False)], spacing=1.35)

# ==========================================================================
# 07 — Full product vision
# ==========================================================================
s = slide()
head(s, "The vision", "What LifeOS is meant to be", 32,
     sub="A single place where planning, deadlines, habits, focus and enforcement live together.")

items = [
    ("Goal expansion", "\u201cCrack Google in a month\u201d becomes tasks, habits, study blocks, caps and alarms.", ACCENT),
    ("Schedule & deadlines", "Every commitment on one timeline, with a live risk score per deadline.", ACCENT),
    ("Mail \u2192 calendar", "Exams and deadlines pulled out of your inbox and added for you.", VIOLET),
    ("Habits & routines", "Bedtime and morning nudges, set dynamically from the day you just had.", VIOLET),
    ("Focus enforcement", "Distracting apps covered by an overlay while a session is live.", DANGER),
    ("Screen time & network", "Daily per-app caps, plus on-device DNS blocking that kills the site too.", DANGER),
    ("Alarms that talk", "Reminders and wake-ups spoken in your chosen persona's voice.", WARN),
    ("Evolving personality", "It shouts or it's kind — and tunes toward whatever actually moves you.", SUCCESS),
    ("Plugins", "WhatsApp, Moodle, Piazza, any site — auto-ingest what you'd otherwise miss.", SUCCESS),
]
gw = (CW - 0.56) / 3
gh = 1.28
for i, (t, b, c) in enumerate(items):
    x = M + (i % 3) * (gw + 0.28)
    y = 2.35 + (i // 3) * (gh + 0.18)
    feature_card(s, x, y, gw, gh, t, b, c, pad=0.22)

# ==========================================================================
# 08 — Personality
# ==========================================================================
s = slide()
head(s, "The relationship", "It shouts, or it's kind.\nYou choose — then it learns.", 30)

pc = (CW - 0.56) / 3
feature_card(s, M, 2.95, pc, 1.95, "Strict",
             "Blunt and terse. Holds you to the deadline. No pep talk, no excuses.", DANGER)
feature_card(s, M + pc + 0.28, 2.95, pc, 1.95, "Supportive",
             "Warm and encouraging. Forgives one slip, then gently gets you back.", SUCCESS)
feature_card(s, M + 2 * (pc + 0.28), 2.95, pc, 1.95, "Coach",
             "Energetic and competitive. Frames the whole week as training.", ACCENT)

card(s, M, 5.25, CW, 1.45, fill=SURFACE2)
text(s, M + 0.35, 5.5, CW - 0.7, 1.0,
     [("The persona is the surface. Underneath, it watches how you respond — "
       "when you fold, what you ignore, how technical you want the answer, how "
       "hard you need to be pushed — and drifts its tone toward what actually "
       "gets you working. Same engine, different voice.", 13, T2, False)], spacing=1.4)

# ==========================================================================
# 09 — Memory / compaction
# ==========================================================================
s = slide()
head(s, "The hard part", "Chat forgets. Your life shouldn't.", 32,
     sub="Every AI wrapper loses your goal the moment its context window is summarized.")

card(s, M, 2.65, 5.35, 2.95)
text(s, M + 0.35, 2.93, 4.6, 0.4, [("EPHEMERAL  ·  COMPACTABLE", 9.5, T3, True)])
bullets(s, M + 0.35, 3.45, 4.6,
        ["Chat transcript, last 20 turns", "Older turns folded into a summary",
         "Raw messages deleted", "Tokens are finite — this must happen"],
        size=11.5, gap=0.48, mcolor=T3, color=T2)

card(s, M + 5.65, 2.65, 5.35, 2.95, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 6.0, 2.93, 4.6, 0.4, [("DURABLE  ·  NEVER COMPACTED", 9.5, ACCENT, True)])
bullets(s, M + 6.0, 3.45, 4.6,
        ["Goals, tasks, events, habits", "Schedule blocks and alarms",
         "App caps and focus rules", "Memory facts about you"],
        size=11.5, gap=0.48, marker="\u2713", mcolor=SUCCESS, color=T2)

text(s, M, 5.95, CW, 0.9,
     [("Your goals live in a database the summarizer cannot touch, and get "
       "re-injected into every single turn. ", 13, T2, False),
      ("Forgetting is a bug we designed out, not a limitation we live with.",
       13, T1, True)], spacing=1.38)

# ==========================================================================
# 10 — Why OS
# ==========================================================================
s = slide()
head(s, "The long game", "Why call it an operating system?", 32)

stages = ["Task manager", "Planner", "Coach", "Accountability partner", "Life companion"]
sw = (CW - 4 * 0.22) / 5
for i, st in enumerate(stages):
    x = M + i * (sw + 0.22)
    c = ACCENT_VIVID if i >= 3 else T3
    card(s, x, 2.75, sw, 0.95, fill=SURFACE2 if i >= 3 else SURFACE1,
         border=ACCENT_DEEP if i >= 3 else BORDER)
    text(s, x + 0.16, 2.95, sw - 0.32, 0.6,
         [(st, 11.5, T1 if i >= 3 else T2, i >= 3)], align=PP_ALIGN.CENTER,
         spacing=1.2)
    if i < 4:
        text(s, x + sw + 0.01, 3.05, 0.22, 0.3, [("\u203a", 14, T3, True)],
             align=PP_ALIGN.CENTER)

card(s, M, 4.15, CW, 2.05, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 0.45, 4.45, CW - 0.9, 1.5,
     [[("Today it runs one life on one phone.", 15, ACCENT, True)],
      [("One day there will be autonomous agents and robots acting on our "
        "behalf. They will need a layer that already knows a person's goals, "
        "constraints, habits and priorities — and is trusted to act on them. "
        "We're building that layer, and starting where the distraction is.",
        13, T2, False)]], spacing=1.42, space_after=8)

# ==========================================================================
# 11 — Retention
# ==========================================================================
s = slide()
head(s, "Why users stay", "Retention isn't a streak.\nIt's a result.", 32)

text(s, M, 2.95, 5.6, 2.6,
     [("When someone actually cracks the interview, or finally loses the 15 kg, "
       "with LifeOS pushing them the whole way — the app stops being a tool "
       "they use.", 14, T1, False)], spacing=1.45)
text(s, M, 4.55, 5.6, 1.0,
     [("It becomes the reason it happened.", 16, ACCENT, True)])

rx = M + 5.95
rw2 = CW - 5.95
card(s, rx, 2.7, rw2, 3.85)
rows = [
    ("Your history is the moat", ACCENT,
     "Weeks of goals, failures, and what actually worked on you. "
     "A fresh chatbot starts at zero."),
    ("Attached to the outcome", SUCCESS,
     "Nobody churns from the thing that got them the offer. "
     "The bond is the result, not the UI."),
    ("The wedge", VIOLET,
     "Students with a hard deadline this month and a documented "
     "distraction problem. High pain, high urgency."),
]
for i, (t, c, b) in enumerate(rows):
    y = 3.05 + i * 1.18
    bar(s, rx + 0.35, y, 0.42, 0.045, c)
    text(s, rx + 0.35, y + 0.19, rw2 - 0.7, 0.35, [(t, 13, T1, True)])
    text(s, rx + 0.35, y + 0.56, rw2 - 0.7, 0.6, [(b, 11, T2, False)], spacing=1.32)

# ==========================================================================
# 12 — Shipped today: the build
# ==========================================================================
s = slide()
head(s, "Built today", "Not a mockup.\nA working Android app.", 32,
     sub="Native Kotlin and Jetpack Compose, nine Gradle modules, installable APK on a real device.")

vals = [("9", "Gradle modules"), ("110", "Kotlin files"), ("26", "typed AI actions"),
        ("6", "app screens"), ("15", "test suites"), ("1", "installable APK")]
sw = (CW - 5 * 0.2) / 6
for i, (v, l) in enumerate(vals):
    stat(s, M + i * (sw + 0.2), 3.0, sw, v, l,
         color=ACCENT if i in (0, 2) else T1)

card(s, M, 4.5, CW, 2.2)
text(s, M + 0.35, 4.78, CW - 0.7, 0.35, [("MODULE GRAPH", 9.5, T3, True)])
mods = [
    ("core", "models, actions, ports"), ("domain", "executor, risk, projection"),
    ("agent", "LLM client, prompts, parser"), ("email", "IMAP, classifier"),
    ("data", "durable store, secrets"), ("enforce", "focus, alarms, VPN"),
    ("calendar", "device calendar sync"), ("ui", "Compose screens"),
    ("app", "wiring, manifest"),
]
for i, (m, d) in enumerate(mods):
    x = M + 0.35 + (i % 3) * ((CW - 0.7) / 3)
    y = 5.25 + (i // 3) * 0.42
    text(s, x, y, (CW - 0.7) / 3 - 0.2, 0.4,
         [(":" + m + "  ", 11.5, ACCENT, True), (d, 10.5, T2, False)], spacing=1.2)

# ==========================================================================
# 13 — Shipped today: what works
# ==========================================================================
s = slide()
head(s, "Working features", "What's live in the APK", 32)

col = (CW - 0.56) / 3
groups = [
    ("AGENT & MEMORY", ACCENT, [
        "Chat with three personas",
        "One sentence \u2192 a whole plan",
        "26 typed actions, tolerant parser",
        "Life state survives chat compaction",
        "Deterministic deadline risk %",
        "Offline fallback if the API dies",
    ]),
    ("ENFORCEMENT", DANGER, [
        "Focus overlay via usage-stats polling",
        "Back button doesn't lift the block",
        "Per-app daily time caps",
        "On-device DNS blocking (VpnService)",
        "Exact alarms, full-screen + spoken line",
        "Survives reboot",
    ]),
    ("INGESTION & SURFACES", SUCCESS, [
        "IMAP fetch + rule-based classifier",
        "Promote an email into your schedule",
        "Google Calendar via device account",
        "Today, Goals, Inbox, Focus, More",
        "Encrypted credential storage",
        "Permission onboarding flow",
    ]),
]
for i, (t, c, its) in enumerate(groups):
    x = M + i * (col + 0.28)
    card(s, x, 2.6, col, 4.1)
    bar(s, x + 0.28, 2.9, 0.42, 0.045, c)
    text(s, x + 0.28, 3.12, col - 0.56, 0.35, [(t, 10, T1, True)])
    for j, it in enumerate(its):
        text(s, x + 0.28, 3.62 + j * 0.5, col - 0.56, 0.5,
             [("\u2713  ", 11, c, False), (it, 10.5, T2, False)], spacing=1.25)

# ==========================================================================
# 14 — Demo
# ==========================================================================
s = slide()
head(s, "See it", "Ninety seconds, six beats", 32)

beats = [
    ("One sentence, twelve changes", "A goal is typed. Goals and Today fill in."),
    ("It's real device state", "Focus on. Leave the app. The overlay slams in front."),
    ("Back doesn't save you", "The block holds; you land on the home screen."),
    ("It kills the site too", "DNS blocked on-device — the browser can't reach it either."),
    ("It wakes you up", "Full-screen alarm, persona line spoken aloud."),
    ("It doesn't forget", "Compact the chat. Message count drops. Life state doesn't."),
]
bh = 0.66
for i, (t, d) in enumerate(beats):
    y = 2.7 + i * (bh + 0.05)
    text(s, M, y + 0.1, 0.5, 0.45, [("%d" % (i + 1), 15, ACCENT_VIVID, True)])
    text(s, M + 0.5, y + 0.06, 4.6, 0.45, [(t, 13, T1, True)])
    text(s, M + 5.3, y + 0.09, CW - 5.3, 0.45, [(d, 12, T2, False)])
    if i < 5:
        bar(s, M + 0.5, y + bh - 0.04, CW - 0.5, 0.008, BORDER)

# ==========================================================================
# 15 — Roadmap + model
# ==========================================================================
s = slide()
head(s, "Next", "From prototype to product", 32)

rw = (CW - 0.56) / 3
feature_card(s, M, 2.7, rw, 2.4, "Now (shipped)",
             "Agent, durable life state, focus overlay, app caps, DNS blocking, "
             "alarms, IMAP inbox, calendar sync.", SUCCESS)
feature_card(s, M + rw + 0.28, 2.7, rw, 2.4, "Next few months",
             "Plugin SDK for WhatsApp, Moodle, Piazza and arbitrary sites. "
             "Adaptive replanning when you fall behind. Personality that drifts "
             "with your history.", ACCENT)
feature_card(s, M + 2 * (rw + 0.28), 2.7, rw, 2.4, "Later",
             "Health signals into planning. Opportunity detection. Chief-of-staff "
             "mode that finds your deadlines before you do.", VIOLET)

card(s, M, 5.35, CW, 1.35, fill=SURFACE2, border=ACCENT_DEEP)
text(s, M + 0.35, 5.6, CW - 0.7, 0.9,
     [("Model: ", 12.5, ACCENT, True),
      ("free core enforcement, paid personality packs and deep integrations. "
       "Focus apps already charge ten dollars a month for a dumb wall — we're "
       "charging for an accountability partner you can't otherwise hire.",
       12.5, T2, False)], spacing=1.4)

# ==========================================================================
# 16 — Close
# ==========================================================================
s = slide(bg=BACKDROP, numbered=False)
bar(s, M, 2.6, 0.055, 1.85)
text(s, M + 0.3, 2.5, 11.0, 1.9,
     [[("ChatGPT writes the plan.", 40, T1, True)],
      [("LifeOS takes the phone.", 40, ACCENT, True)]], spacing=1.16)
text(s, M + 0.32, 4.8, 9.0, 0.5,
     [("Stop managing tasks. Start executing goals.", 16, T2, False)])
text(s, M, 6.35, 9.0, 0.4,
     [("Working APK, on a real device, right now.", 11, T3, False)])

# ==========================================================================
# Speaker notes
# ==========================================================================
NOTES = [
    "[0:00] Don't read the slide. Say the name, say the tagline, move on. "
    "Energy here sets the whole pitch.",

    "[0:10] Ask it as a real question and pause for the nod. This is a personal "
    "story, not a market stat — tell it that way. Land 'At least, I do' and let "
    "it sit for a beat.",

    "[0:25] The reframe. Everyone in the room has tried a productivity app and "
    "quit. Name that. The point: organization is solved, execution isn't. "
    "Say 'Nothing' out loud and let it hang.",

    "[0:40] This is why the gap still exists in 2026 — it's structural, not a "
    "gap nobody noticed. Incumbents are financially incapable of shipping this. "
    "That's our permission to exist.",

    "[0:55] The core demo idea in one slide. Read the sentence aloud as if you "
    "were typing it. Then hit the right column hard: goal, tasks, cap, alarm, "
    "memory. Close on 'it reaches into the phone and changes it.'",

    "[1:15] Only slide with architecture. Keep it to two sentences: one model "
    "call returns a reply plus typed actions, and exactly one component applies "
    "them. The 'database is the soul' line is the takeaway.",

    "[1:30] Breadth slide — do NOT read all nine. Sweep a hand across it and "
    "name three: mail to calendar, screen time and network, plugins. Say the "
    "rest is on the slide.",

    "[1:45] The fun one. Mention the personas by name. Then the real point: it "
    "drifts toward whatever actually moves you. This is where retention starts.",

    "[2:00] Technical credibility. Every judge has used an AI that forgot "
    "something. Explain: chat gets summarized, life state never does, and it's "
    "re-injected every turn. This is a real engineering choice, not a slogan.",

    "[2:15] Answers 'why the grand name'. Say the progression, then the honest "
    "long game: agents will need a layer that knows a person's goals. Don't "
    "oversell — 'hopefully' is fine, confidence is better.",

    "[2:30] Business credibility. The insight: retention comes from the "
    "outcome, not gamification. Someone who got the offer with LifeOS never "
    "churns. Name the wedge so it doesn't sound like 'everyone'.",

    "[2:45] Pivot from vision to proof. Say plainly: this isn't slideware, "
    "there's an APK. Numbers are real — nine modules, 110 files, 26 action "
    "types, six screens, tests.",

    "[3:00] Don't read all eighteen. Point at each column header and give one "
    "example. Highlight the two nobody expects: the back button doesn't lift "
    "the block, and DNS blocking works in the browser too.",

    "[3:15] Transition into the live demo. If demoing, stop talking and show "
    "beats 1-3 minimum. If time is short, show the overlay beat — that's the "
    "one people remember.",

    "[3:30] Forward-looking. Now / next / later. Then the money line: focus "
    "apps already charge for a dumb wall; we're charging for an accountability "
    "partner you can't hire.",

    "[3:45] Land it and stop. Say both lines, then the tagline, then be quiet. "
    "Don't add anything after 'LifeOS takes the phone.'",
]
for _sl, _nt in zip(prs.slides, NOTES):
    _sl.notes_slide.notes_text_frame.text = _nt

prs.save("/home/sumit/lifeos/pitch/LifeOS_Pitch.pptx")
print("wrote pitch/LifeOS_Pitch.pptx  (%d slides)" % len(prs.slides.__iter__.__self__._sldIdLst))
